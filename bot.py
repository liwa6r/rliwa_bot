"""
بوت ليوا v3 - نظام تقارير التفتيش والتوعية
"""

import os
import re
import random
import base64
import logging
import smtplib
import asyncio
import requests as req
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email.mime.text import MIMEText
from email import encoders
from datetime import datetime
from io import BytesIO

import anthropic
from telegram import Update, ReplyKeyboardMarkup, KeyboardButton, ReplyKeyboardRemove
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    filters, ContextTypes, ConversationHandler
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage

logging.basicConfig(format="%(asctime)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)

# ── الإصدار
BOT_VERSION = "v3.4"
BOT_NAME    = "بوت ليوا"

# ── أكواد الدخول
CODE_ADMIN      = "9090"   # مدير النظام — فحص شامل
CODE_AWARENESS  = "9091"   # قسم التوعية
CODE_INSPECTION = "9092"   # قسم التفتيش

# ── المتغيرات البيئية
TELEGRAM_TOKEN    = os.environ.get("TELEGRAM_TOKEN", "")
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
ADMIN_EMAIL       = "bot.6rr@gmail.com"
GMAIL_USER        = os.environ.get("GMAIL_USER", "")
GMAIL_APP_PASS    = os.environ.get("GMAIL_APP_PASS", "")
RESEND_API_KEY    = os.environ.get("RESEND_API_KEY", "")
GOOGLE_MAPS_KEY   = os.environ.get("GOOGLE_MAPS_KEY", "")

# ── نطاق ليوا والمنطقة الغربية
LIWA_BOUNDS = {
    "lat_min": 22.0, "lat_max": 24.5,
    "lng_min": 51.5, "lng_max": 55.5
}

# ── ألوان النموذج
BROWN    = RGBColor(0x8B, 0x25, 0x00)
LIGHT_BG = RGBColor(0xF2, 0xF0, 0xED)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLACK    = RGBColor(0x1A, 0x1A, 0x1A)
TEAL     = RGBColor(0x00, 0x6B, 0x6B)   # لون تقارير التوعية

# ── مراحل المحادثة
(WAITING_CODE, WAITING_NAME, MAIN_MENU,
 WAITING_PHOTO, WAITING_LOCATION, WAITING_AWARENESS_TYPE,
 WAITING_NOTE, WAITING_PLOT, ADD_MORE,
 WAITING_BEFORE_PHOTO, WAITING_AFTER_PHOTO,
 WAITING_BA_LOC, WAITING_BA_NOTE,
 CONFIRM_SEND, COLLECTING_PHOTOS,
 WAITING_EMAIL) = range(16)

user_sessions = {}
report_stats  = {"total": 0, "inspection": 0, "awareness": 0}


# ══════════════════════════════════════════════════
#  Google Geocoding
# ══════════════════════════════════════════════════

def decimal_to_dms(lat: float, lng: float) -> str:
    """تحويل الإحداثيات العشرية إلى صيغة DMS"""
    def to_dms(deg, pos, neg):
        d = int(abs(deg))
        m = int((abs(deg) - d) * 60)
        s = round(((abs(deg) - d) * 60 - m) * 60, 1)
        hemi = pos if deg >= 0 else neg
        return f"{d}\u00b0{m}'{s}\"\u202a{hemi}\u202c"
    return f"{to_dms(lat, 'N', 'S')} {to_dms(lng, 'E', 'W')}"


def get_address(lat: float, lng: float) -> dict:
    in_liwa = (
        LIWA_BOUNDS["lat_min"] <= lat <= LIWA_BOUNDS["lat_max"] and
        LIWA_BOUNDS["lng_min"] <= lng <= LIWA_BOUNDS["lng_max"]
    )
    location_name = f"{lat:.4f}, {lng:.4f}"

    if GOOGLE_MAPS_KEY:
        try:
            url = (
                f"https://maps.googleapis.com/maps/api/geocode/json"
                f"?latlng={lat},{lng}&language=ar&key={GOOGLE_MAPS_KEY}"
            )
            r = req.get(url, timeout=6)
            data = r.json()
            if data.get("status") == "OK" and data.get("results"):
                components = data["results"][0].get("address_components", [])
                route = sublocality = locality = ""
                for c in components:
                    types = c.get("types", [])
                    name  = c.get("long_name", "")
                    if "route" in types and not route:
                        route = name
                    elif ("sublocality" in types or "sublocality_level_1" in types) and not sublocality:
                        sublocality = name
                    elif "locality" in types and not locality:
                        locality = name
                parts = [p for p in [route, sublocality] if p]
                if not parts and locality:
                    parts = [locality]
                if parts:
                    location_name = "، ".join(parts)
                else:
                    full = data["results"][0].get("formatted_address", "")
                    full = re.sub(r"الإمارات.*|Abu Dhabi.*|UAE.*", "", full)
                    location_name = full.strip().strip(",").strip()
        except Exception as e:
            logger.error(f"Geocoding: {e}")

    return {"name": location_name, "in_liwa": in_liwa, "lat": lat, "lng": lng}


def parse_coords(text: str):
    # أولاً: توسيع الروابط المختصرة
    if "maps.app.goo.gl" in text or "goo.gl" in text:
        try:
            r = req.get(text.strip(), allow_redirects=True, timeout=8)
            text = r.url  # الرابط الكامل بعد التوجيه
        except Exception:
            pass

    patterns = [
        r'(-?\d+\.?\d*)[,،\s]+(-?\d+\.?\d*)',
        r'[?&]q=(-?\d+\.?\d*),(-?\d+\.?\d*)',
        r'@(-?\d+\.?\d*),(-?\d+\.?\d*)',
        r'll=(-?\d+\.?\d*),(-?\d+\.?\d*)',
        r'!3d(-?\d+\.?\d*)!4d(-?\d+\.?\d*)',
    ]
    for p in patterns:
        m = re.search(p, text)
        if m:
            la, ln = float(m.group(1)), float(m.group(2))
            if -90 <= la <= 90 and -180 <= ln <= 180:
                return la, ln
    return None, None


# ══════════════════════════════════════════════════
#  Claude AI
# ══════════════════════════════════════════════════

def ai_reply(msg: str, session: dict) -> str:
    try:
        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        dept   = "التوعية" if session.get("dept") == "awareness" else "التفتيش"
        sys = (
            f"أنت {BOT_NAME} {BOT_VERSION} — مساعد تقارير {dept} في مدينة ليوا.\n"
            f"إجمالي التقارير: {report_stats['total']} "
            f"(تفتيش: {report_stats['inspection']} — توعية: {report_stats['awareness']})\n"
            "أجب بالعربية بشكل ودي ومختصر (جملتان فقط)."
        )
        r = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=150,
            system=sys,
            messages=[{"role": "user", "content": msg}]
        )
        return r.content[0].text.strip()
    except Exception:
        return f"أنا {BOT_NAME} {BOT_VERSION} — كيف أساعدك؟"


# ══════════════════════════════════════════════════
#  Agent الذكي — مراجعة وضبط بيانات التقرير
# ══════════════════════════════════════════════════

MAX_IMG_BYTES = 3 * 1024 * 1024   # 3 MB حد أقصى لحجم الصورة
MAX_IMG_DIM   = 1920               # حد أقصى للبعد بالبكسل


def agent_fix_image(photo_bytes: bytes) -> bytes:
    """ضبط حجم وأبعاد الصورة لتناسب التقرير"""
    try:
        img = PILImage.open(BytesIO(photo_bytes))
        # تحويل لـ RGB لو كان RGBA أو غيره
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        # تصغير لو الأبعاد كبيرة جداً
        w, h = img.size
        if max(w, h) > MAX_IMG_DIM:
            ratio = MAX_IMG_DIM / max(w, h)
            img = img.resize((int(w*ratio), int(h*ratio)), PILImage.LANCZOS)
        # ضغط لو الحجم كبير
        out = BytesIO()
        quality = 85
        img.save(out, format="JPEG", quality=quality, optimize=True)
        # لو لازال كبير اخفض الجودة أكثر
        while out.tell() > MAX_IMG_BYTES and quality > 50:
            quality -= 10
            out = BytesIO()
            img.save(out, format="JPEG", quality=quality, optimize=True)
        out.seek(0)
        return out.read()
    except Exception as e:
        logger.warning(f"agent_fix_image: {e}")
        return photo_bytes  # أرجع الأصلي لو فشل


def agent_fix_text(text: str) -> str:
    """تصحيح الإملاء وتحسين الصياغة عبر Claude"""
    if not text or not text.strip() or not ANTHROPIC_API_KEY:
        return text
    try:
        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        r = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=200,
            system=(
                "أنت مدقق لغوي متخصص في تقارير التفتيش الميداني. "
                "صحح الإملاء والتشكيل فقط دون تغيير المعنى. "
                "أرجع النص المصحح فقط بدون شرح أو إضافات."
            ),
            messages=[{"role": "user", "content": text}]
        )
        fixed = r.content[0].text.strip()
        # تأكد إن الرد نص قصير وليس شرحاً طويلاً
        if fixed and len(fixed) < len(text) * 3:
            return fixed
        return text
    except Exception as e:
        logger.warning(f"agent_fix_text: {e}")
        return text


def agent_prepare_session(session: dict) -> dict:
    """
    Agent رئيسي: يمر على كل صورة في الجلسة ويضبط:
    1. حجم الصورة وأبعادها
    2. إملاء الملاحظة
    3. إملاء نوع التوعية
    """
    import copy
    prepared = copy.deepcopy(session)
    for entry in prepared.get("photos", []):
        # ضبط الصورة الرئيسية
        if entry.get("photo"):
            entry["photo"] = agent_fix_image(entry["photo"])
        # ضبط صورة البعد (قبل/بعد)
        if entry.get("photo_after"):
            entry["photo_after"] = agent_fix_image(entry["photo_after"])
        # تصحيح الملاحظة
        if entry.get("note"):
            entry["note"] = agent_fix_text(entry["note"])
        # تصحيح نوع التوعية
        if entry.get("awareness_type"):
            entry["awareness_type"] = agent_fix_text(entry["awareness_type"])
    return prepared


# ══════════════════════════════════════════════════
#  فحص النظام — كود المدير 9090
# ══════════════════════════════════════════════════

async def run_system_check(update: Update):
    await update.message.reply_text("🔍 جاري فحص النظام... لحظة")

    results = []
    all_ok  = True

    # 1. Claude AI
    try:
        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
        client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=5,
            messages=[{"role": "user", "content": "ping"}]
        )
        results.append("✅ Claude AI — يعمل")
    except Exception as e:
        results.append(f"❌ Claude AI — {str(e)[:50]}")
        all_ok = False

    # 2. Google Maps
    try:
        url = (
            f"https://maps.googleapis.com/maps/api/geocode/json"
            f"?latlng=23.085,54.016&language=ar&key={GOOGLE_MAPS_KEY}"
        )
        r    = req.get(url, timeout=6)
        data = r.json()
        if data.get("status") == "OK":
            addr = data["results"][0].get("formatted_address", "")
            addr = re.sub(r"الإمارات.*|UAE.*", "", addr).strip()[:45]
            results.append(f"✅ Google Maps — يعمل\n    📍 {addr}")
        else:
            results.append(f"❌ Google Maps — {data.get('status')} (تحقق من Billing)")
            all_ok = False
    except Exception as e:
        results.append(f"❌ Google Maps — {str(e)[:50]}")
        all_ok = False

    # 3. Resend API
    try:
        r = req.get(
            "https://api.resend.com/domains",
            headers={"Authorization": f"Bearer {RESEND_API_KEY}"},
            timeout=8
        )
        if r.status_code == 200:
            results.append(f"✅ Resend — متصل ✓")
        else:
            results.append(f"❌ Resend — خطأ {r.status_code}")
            all_ok = False
    except Exception as e:
        results.append(f"❌ Resend — {str(e)[:50]}")
        all_ok = False

    # 4. إيميل تجريبي حقيقي
    if RESEND_API_KEY:
        try:
            r = req.post(
                "https://api.resend.com/emails",
                headers={
                    "Authorization": f"Bearer {RESEND_API_KEY}",
                    "Content-Type": "application/json"
                },
                json={
                    "from":    "بوت ليوا <onboarding@resend.dev>",
                    "to":      [ADMIN_EMAIL],
                    "subject": f"🔧 فحص النظام — {BOT_NAME} {BOT_VERSION}",
                    "text": (
                        f"تقرير فحص النظام\n\n"
                        f"البوت: {BOT_NAME} {BOT_VERSION}\n"
                        f"التاريخ: {datetime.now().strftime('%Y/%m/%d %H:%M')}\n\n"
                        f"إجمالي التقارير: {report_stats['total']}\n"
                        f"تقارير التفتيش: {report_stats['inspection']}\n"
                        f"تقارير التوعية: {report_stats['awareness']}\n\n"
                        + "\n".join(results)
                    )
                },
                timeout=15
            )
            if r.status_code == 200:
                results.append(f"✅ إيميل تجريبي — وصل إلى {ADMIN_EMAIL}")
            else:
                results.append(f"❌ إيميل تجريبي — {r.text[:50]}")
                all_ok = False
        except Exception as e:
            results.append(f"❌ إيميل تجريبي — {str(e)[:50]}")
            all_ok = False

    status = "🟢 كل شيء يعمل بشكل صحيح" if all_ok else "🔴 يوجد مشاكل تحتاج مراجعة"

    report = (
        f"━━━━━━━━━━━━━━━━━━━━\n"
        f"🔧 {BOT_NAME} {BOT_VERSION}\n"
        f"📊 تقرير صحة النظام\n"
        f"━━━━━━━━━━━━━━━━━━━━\n\n"
        + "\n".join(results) +
        f"\n\n━━━━━━━━━━━━━━━━━━━━\n"
        f"📈 إجمالي التقارير: {report_stats['total']}\n"
        f"🔍 تفتيش: {report_stats['inspection']}  |  📢 توعية: {report_stats['awareness']}\n"
        f"🕐 {datetime.now().strftime('%Y/%m/%d — %H:%M')}\n\n"
        f"{status}"
    )
    await update.message.reply_text(report)


# ══════════════════════════════════════════════════
#  نظام الأوامر الذكي
# ══════════════════════════════════════════════════

# كلمات استفتاحية تبدأ البوت تلقائياً
START_KEYWORDS = [
    "مرحبا", "مرحباً", "هلا", "السلام", "ابدأ", "ابداء", "ابدا",
    "تقرير", "توعية", "تفتيش", "بدء", "انطلق", "يلا", "هيا",
    "start", "بداية", "جديد", "اهلا", "أهلاً", "أهلا", ".", "،"
]

# أوامر المسح والإعادة
RESET_KEYWORDS = [
    "مسح", "إعادة", "اعادة", "من جديد", "مسح التقرير",
    "عيد من اول", "أعد من أول", "اعد من اول", "إعادة تشغيل",
    "حذف", "صفي", "امسح", "امسح كل شي", "حذف كل شي"
]

# أوامر تصحيح الأخطاء
ERROR_KEYWORDS = ["هناك خطأ", "هناك خطاء", "في خطأ", "خطأ", "خطاء", "غلط"]
ERROR_PHOTO    = ["خطأ في الصورة", "خطاء في الصورة", "غلط الصورة", "صورة غلط", "بدل الصورة"]
ERROR_LOCATION = ["خطأ في الموقع", "خطاء في الموقع", "غلط الموقع", "موقع غلط", "بدل الموقع"]
ERROR_AWARENESS= ["خطأ في التوعية", "خطاء في التوعية", "غلط التوعية", "بدل التوعية"]

# أوامر الحالة
STATUS_KEYWORDS   = ["الحالة", "حالة", "وين وصلت", "كم صورة", "كم الصور"]
DELETE_LAST       = ["احذف آخر صورة", "احذف الاخيرة", "امسح آخر صورة", "احذف اخر صورة"]
PREVIEW_KEYWORDS  = ["معاينة", "معاينه", "اعرض التقرير", "شوف التقرير"]
HELP_KEYWORDS     = ["مساعدة", "مساعده", "help", "الأوامر", "الاوامر", "ايش الاوامر"]

# أوامر المدير
ADMIN_STATS       = ["إحصائيات", "احصائيات", "تقارير اليوم", "الإحصائيات"]


def match_keywords(text: str, keywords: list) -> bool:
    text_lower = text.lower().strip()
    return any(kw in text_lower for kw in keywords)


async def handle_commands(update: Update, context: ContextTypes.DEFAULT_TYPE, current_state: int):
    """يعالج الأوامر الذكية في أي مرحلة"""
    user_id = update.effective_user.id
    text    = update.message.text.strip()
    session = user_sessions.get(user_id, {})

    # ── كلمات استفتاحية — ابدأ البوت
    if match_keywords(text, START_KEYWORDS) and not session:
        await start(update, context)
        return WAITING_CODE

    # ── مسح وإعادة
    if match_keywords(text, RESET_KEYWORDS):
        user_sessions.pop(user_id, None)
        await update.message.reply_text(
            "🌴 مرحباً بك في بوت مدينة ليوا\n\n"
            "🔐 اكتب الكود الخاص بك للدخول:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_CODE

    # ── مساعدة
    if match_keywords(text, HELP_KEYWORDS):
        dept = "التوعية" if session.get("dept") == "awareness" else "التفتيش"
        await update.message.reply_text(
            f"📋 *قائمة الأوامر — {BOT_NAME} {BOT_VERSION}*\n\n"
            "━━━━━━━━━━━━━━━━━━━\n"
            "🔄 *إعادة وتصحيح:*\n"
            "• `مسح` — يمسح ويبدأ من أول\n"
            "• `هناك خطأ` — يسألك عن نوع الخطأ\n"
            "• `خطأ في الصورة` — يرجع لخطوة الصورة\n"
            "• `خطأ في الموقع` — يرجع لخطوة الموقع\n"
            "• `خطأ في التوعية` — يرجع لخطوة التوعية\n"
            "• `احذف آخر صورة` — يحذف الصورة الأخيرة\n\n"
            "📊 *معلومات:*\n"
            "• `الحالة` — ملخص التقرير الحالي\n"
            "• `معاينة` — معاينة التقرير\n"
            "• `كم صورة` — عدد الصور المضافة\n\n"
            "━━━━━━━━━━━━━━━━━━━\n"
            f"القسم الحالي: {dept}",
            parse_mode="Markdown"
        )
        return current_state

    # ── الحالة
    if match_keywords(text, STATUS_KEYWORDS):
        photos = session.get("photos", [])
        if not photos:
            await update.message.reply_text("📋 لا توجد صور مضافة بعد.")
        else:
            locs = list(dict.fromkeys([p["location"]["name"] for p in photos if p.get("location")]))
            status_msg = (
                f"📊 *حالة التقرير الحالي:*\n\n"
                f"📸 الصور: {len(photos)}\n"
                f"📍 المواقع: {', '.join(locs) or '—'}\n"
                f"📅 التاريخ: {session.get('date','—')}\n"
                f"👤 الاسم: {session.get('name','—')}\n"
            )
            if session.get("dept") == "awareness":
                types = [p.get("awareness_type","") for p in photos if p.get("awareness_type")]
                status_msg += f"📋 أنواع التوعية: {', '.join(types) or '—'}\n"
            await update.message.reply_text(status_msg, parse_mode="Markdown")
        return current_state

    # ── معاينة
    if match_keywords(text, PREVIEW_KEYWORDS):
        photos = session.get("photos", [])
        if not photos:
            await update.message.reply_text("❌ لا توجد صور للمعاينة.")
        else:
            preview = f"👁️ *معاينة التقرير:*\n\n"
            for i, p in enumerate(photos, 1):
                preview += f"📸 *صورة {i}:*\n"
                preview += f"   📍 {p['location']['name']}\n"
                if p.get("awareness_type"):
                    preview += f"   📋 {p['awareness_type']}\n"
                if p.get("note"):
                    preview += f"   📝 {p['note']}\n"
                preview += "\n"
            await update.message.reply_text(preview, parse_mode="Markdown")
        return current_state

    # ── احذف آخر صورة
    if match_keywords(text, DELETE_LAST):
        photos = session.get("photos", [])
        if not photos:
            await update.message.reply_text("❌ لا توجد صور للحذف.")
        else:
            photos.pop()
            await update.message.reply_text(
                f"🗑️ تم حذف الصورة الأخيرة\n"
                f"📸 الصور المتبقية: {len(photos)}"
            )
        return current_state

    # ── هناك خطأ
    if match_keywords(text, ERROR_KEYWORDS):
        kb = [
            [KeyboardButton("📸 خطأ في الصورة")],
            [KeyboardButton("📍 خطأ في الموقع")],
            [KeyboardButton("📋 خطأ في التوعية")],
        ]
        await update.message.reply_text(
            "⚠️ أين الخطأ؟",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return current_state

    # ── تصحيح الصورة
    if match_keywords(text, ERROR_PHOTO):
        # احذف آخر صورة وارجع لخطوة الصورة
        photos = session.get("photos", [])
        if photos:
            photos.pop()
        await update.message.reply_text(
            "📸 أرسل الصورة الصحيحة:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_PHOTO

    # ── تصحيح الموقع
    if match_keywords(text, ERROR_LOCATION):
        photos = session.get("photos", [])
        if photos:
            # احتفظ بالصورة لكن امسح الموقع
            last = photos[-1]
            session.setdefault("collected_photos", []).insert(0, last["photo"])
            session["collecting"] = True
            photos.pop()
        await ask_location(update)
        return WAITING_LOCATION

    # ── تصحيح التوعية
    if match_keywords(text, ERROR_AWARENESS):
        await update.message.reply_text(
            "📋 أدخل نوع التوعية الصحيح:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_AWARENESS_TYPE

    # ── إحصائيات المدير
    if match_keywords(text, ADMIN_STATS):
        await update.message.reply_text(
            f"📊 *الإحصائيات:*\n\n"
            f"📈 إجمالي التقارير: {report_stats['total']}\n"
            f"🔍 تفتيش: {report_stats['inspection']}\n"
            f"📢 توعية: {report_stats['awareness']}\n"
            f"🕐 {datetime.now().strftime('%Y/%m/%d — %H:%M')}",
            parse_mode="Markdown"
        )
        return current_state

    return None  # لم يتطابق أي أمر

def new_session(user_id, name, dept):
    user_sessions[user_id] = {
        "photos":           [],
        "name":             name,
        "dept":             dept,
        "report_type":      "normal",
        "date":             datetime.now().strftime("%Y/%m/%d"),
        "last_location":    None,
        "last_coords_dms":  "",
        "pending_plot":     "",
        "collecting":          False,
        "collected_photos":    [],
        "collecting_msg_id":   None,
        "last_batch_count":    1,
        "active_group_id":  None,
        "group_msg_id":     None,
    }


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        f"🌴 مرحباً بك عزيزي في بوت مدينة ليوا\n\n"
        "🔐 اكتب الكود الخاص بك للدخول:",
        reply_markup=ReplyKeyboardRemove()
    )
    return WAITING_CODE


async def receive_code(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    code    = update.message.text.strip()
    name    = update.effective_user.first_name or "مستخدم"

    # تحقق من الأوامر أولاً
    cmd = await handle_commands(update, context, WAITING_CODE)
    if cmd is not None:
        return cmd

    # ── كلمات استفتاحية
    if match_keywords(code, START_KEYWORDS):
        await update.message.reply_text(
            "🌴 مرحباً بك عزيزي في بوت مدينة ليوا\n\n"
            "🔐 اكتب الكود الخاص بك للدخول:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_CODE

    # ── مدير النظام
    if code == CODE_ADMIN:
        await update.message.reply_text(
            f"🔓 مرحباً بك مدير النظام!\n{BOT_NAME} {BOT_VERSION}",
            reply_markup=ReplyKeyboardRemove()
        )
        await run_system_check(update)
        return ConversationHandler.END

    # ── قسم التفتيش
    elif code == CODE_INSPECTION:
        # نحفظ القسم مؤقتاً حتى نستلم الاسم
        user_sessions[user_id] = {"dept": "inspection", "pending_name": True}
        await update.message.reply_text(
            "✅ كود صحيح! قسم الرصد\n✏️ أدخل اسمك:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_NAME

    # ── قسم التوعية
    elif code == CODE_AWARENESS:
        user_sessions[user_id] = {"dept": "awareness", "pending_name": True}
        await update.message.reply_text(
            "✅ كود صحيح! قسم التوعية\n✏️ أدخل اسمك:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_NAME

    # ── كود خاطئ
    else:
        await update.message.reply_text(
            "❌ الكود غير صحيح\nأعد المحاولة أو تواصل مع المسؤول"
        )
        return WAITING_CODE


# ══════════════════════════════════════════════════
#  استقبال الاسم
# ══════════════════════════════════════════════════

async def receive_name(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})
    name    = update.message.text.strip()

    if not name or len(name) < 2:
        await update.message.reply_text("❌ الاسم قصير جداً. أدخل اسمك كاملاً:")
        return WAITING_NAME

    dept = session.get("dept", "inspection")
    new_session(user_id, name, dept)

    if dept == "inspection":
        kb = [
            [KeyboardButton("📸 تقرير عادي")],
            [KeyboardButton("🔄 تقرير قبل وبعد")],
        ]
        await update.message.reply_text(
            f"أهلاً {name}! اختر نوع التقرير:",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return MAIN_MENU
    else:
        kb = [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]]
        await update.message.reply_text(
            f"أهلاً {name}! 📍\n\nأرسل الموقع أولاً ثم أرسل الصور:\n"
            "1️⃣ اضغط مشاركة موقعي\n"
            "2️⃣ أو رابط جوجل ماب\n"
            "3️⃣ أو إحداثيات: `23.085, 54.016`\n\n"
            "⚠️ يجب أن يكون في نطاق مدينة ليوا",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_LOCATION


# ══════════════════════════════════════════════════
#  القائمة الرئيسية (تفتيش فقط)
# ══════════════════════════════════════════════════

async def main_menu(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text    = update.message.text
    session = user_sessions.get(user_id, {})

    cmd = await handle_commands(update, context, MAIN_MENU)
    if cmd is not None:
        return cmd

    if "عادي" in text or "📸" in text:
        session["report_type"] = "normal"
        kb = [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]]
        await update.message.reply_text(
            "📍 أرسل الموقع أولاً ثم أرسل الصور:\n"
            "1️⃣ اضغط مشاركة موقعي\n"
            "2️⃣ أو رابط جوجل ماب\n"
            "3️⃣ أو إحداثيات: `23.085, 54.016`\n\n"
            "⚠️ يجب أن يكون في نطاق مدينة ليوا",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_LOCATION
    elif "قبل وبعد" in text or "🔄" in text:
        session["report_type"] = "before_after"
        kb = [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]]
        await update.message.reply_text(
            "📍 تقرير قبل وبعد الرصد\n\n"
            "أرسل الموقع أولاً ثم أرسل صور *القبل*:\n"
            "1️⃣ اضغط مشاركة موقعي\n"
            "2️⃣ أو رابط جوجل ماب\n"
            "3️⃣ أو إحداثيات: `23.085, 54.016`",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_LOCATION
    else:
        reply = ai_reply(text, session)
        kb = [[KeyboardButton("📸 تقرير عادي")], [KeyboardButton("🔄 تقرير قبل وبعد")]]
        await update.message.reply_text(reply, reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True))
        return MAIN_MENU


# ══════════════════════════════════════════════════
#  استقبال الصور
# ══════════════════════════════════════════════════
#  معالج media group (صور متعددة دفعة واحدة)
# ══════════════════════════════════════════════════
def _collect_pending_photos(session: dict) -> list:
    """يُرجع الصور المجمّعة"""
    return list(session.get("collected_photos", []))


def _build_group_keyboard(count: int, has_location: bool, loc_name: str = "") -> tuple:
    """يبني رسالة وكيبورد استلام الصور المتعددة"""
    if has_location:
        kb = [
            [KeyboardButton(f"✅ نفس الموقع ({count} صور)")],
            [KeyboardButton("📍 موقع جديد")]
        ]
        text = (
            f"📸 *استلمت {count} صور* — هل أرسلت الكل؟\n"
            f"📍 الموقع السابق: *{loc_name}*\n\n"
            "اختر:"
        )
    else:
        kb = [
            [KeyboardButton("✅ انتهيت — أرسل الموقع الآن")],
            [KeyboardButton("📍 مشاركة موقعي", request_location=True)]
        ]
        text = (
            f"📸 *استلمت {count} صور* — هل أرسلت الكل؟\n\n"
            "عندما تنتهي من إرسال الصور، اضغط:\n"
            "• *انتهيت* ثم أرسل موقعك\n"
            "• أو شارك موقعك مباشرة"
        )
    return text, ReplyKeyboardMarkup(kb, resize_keyboard=True)

async def receive_photo(update, context):
    """استلام الصور — في وضع الجمع تُضاف مباشرة، وإلا تُطلب الموقع أولاً"""
    user_id = update.effective_user.id
    session = user_sessions.get(user_id)
    if not session:
        await start(update, context)
        return WAITING_CODE

    # تنزيل الصورة
    photo = update.message.photo[-1]
    file  = await context.bot.get_file(photo.file_id)
    pb    = await file.download_as_bytearray()
    photo_bytes = bytes(pb)

    # ── وضع الجمع: الموقع محدد مسبقاً — أضف الصورة مباشرة
    if session.get("collecting"):
        session["collected_photos"].append(photo_bytes)
        count = len(session["collected_photos"])

        # احذف رسالة العداد القديمة وأرسل جديدة بالعدد المحدّث
        old_msg = session.get("collecting_msg_id")
        if old_msg:
            try:
                await context.bot.delete_message(
                    chat_id=update.effective_chat.id,
                    message_id=old_msg
                )
            except Exception:
                pass

        btn_label = "✅ انتهيت من صور القبل" if session.get("report_type") == "before_after" else "✅ انتهيت من الصور"
        sent = await context.bot.send_message(
            chat_id=update.effective_chat.id,
            text=f"📸 استلمت *{count}* صورة\nأرسل المزيد أو اضغط الزر أدناه 👇",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(
                [[KeyboardButton(btn_label)]],
                resize_keyboard=True,
                one_time_keyboard=False
            )
        )
        session["collecting_msg_id"] = sent.message_id
        return COLLECTING_PHOTOS

    # ── لا يوجد موقع بعد — اطلب الموقع أولاً
    await update.message.reply_text(
        "⚠️ أرسل الموقع أولاً ثم أرسل الصور:",
        reply_markup=ReplyKeyboardMarkup(
            [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]],
            resize_keyboard=True
        )
    )
    return WAITING_LOCATION


async def ask_location(update):
    kb = [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]]
    await update.message.reply_text(
        "📍 *الموقع إجباري* — أرسله بإحدى الطرق:\n\n"
        "1️⃣ اضغط 'مشاركة موقعي'\n"
        "2️⃣ أرسل رابط جوجل ماب\n"
        "3️⃣ أرسل إحداثيات: `23.085, 54.016`\n\n"
        "⚠️ يجب أن يكون في نطاق مدينة ليوا",
        parse_mode="Markdown",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )


async def receive_location_gps(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions.get(user_id)
    if not session:
        return await start(update, context)
    loc = update.message.location
    ok = await process_location(update, session, loc.latitude, loc.longitude)
    if ok:
        return COLLECTING_PHOTOS   # انتظر الصور أولاً
    return WAITING_LOCATION


async def receive_location_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions.get(user_id)
    if not session:
        return await start(update, context)

    text = update.message.text.strip()

    if "نفس الموقع" in text and session.get("last_location"):
        loc      = session["last_location"]
        coords   = session.get("last_coords_dms", "")
        all_photos = _collect_pending_photos(session)
        for i, ph in enumerate(all_photos):
            session["photos"].append({
                "photo":          ph,
                "location":       loc,
                "coords_dms":     coords,
                "plot_no":        session.get("pending_plot", "") if i == 0 else "",
                "note":           "",
                "awareness_type": "",
                "type":           "normal"
            })
        session["collected_photos"] = []
        session["collecting"]       = False
        session["pending_plot"]     = ""
        session["active_group_id"]  = None
        session["group_msg_id"]     = None
        return await next_after_location(session, update)

    if "موقع جديد" in text:
        await ask_location(update)
        return WAITING_LOCATION

    # ── انتهيت من الصور — اطلب الموقع
    if "انتهيت" in text or "انتهيت من صور القبل" in text or "انتهيت من الصور" in text:
        if not session.get("collected_photos"):
            await update.message.reply_text("❌ لا توجد صور! أرسل الصور أولاً.")
            return WAITING_PHOTO
        await ask_location(update)
        return WAITING_LOCATION

    lat, lng = parse_coords(text)
    if lat:
        ok = await process_location(update, session, lat, lng)
        if ok:
            return COLLECTING_PHOTOS   # انتظر الصور أولاً
        return WAITING_LOCATION

    await update.message.reply_text(
        "❌ لم أتعرف على الموقع.\nجرب رابط جوجل ماب أو إحداثيات مثل:\n`23.085, 54.016`",
        parse_mode="Markdown"
    )
    return WAITING_LOCATION


async def process_location(update, session, lat, lng):
    await update.message.reply_text("⏳ جاري تحديد الموقع...")
    loc_data = get_address(lat, lng)

    if not loc_data["in_liwa"]:
        await update.message.reply_text(
            f"⚠️ *الموقع خارج نطاق ليوا!*\n`{lat:.4f}, {lng:.4f}`\n\nتأكد من موقعك وأعد الإرسال.",
            parse_mode="Markdown"
        )
        await ask_location(update)
        return False

    coords = decimal_to_dms(lat, lng)
    session["last_location"]     = loc_data
    session["last_coords_dms"]   = coords
    session["collecting"]        = True
    session["collected_photos"]  = []
    session["collecting_msg_id"] = None

    if session.get("report_type") == "before_after":
        await update.message.reply_text(
            f"✅ *{loc_data['name']}*\n\n"
            "📸 أرسل صور *القبل* — واحدة أو أكثر\n"
            "⬇️ عند الانتهاء اضغط الزر أدناه:",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(
                [[KeyboardButton("✅ انتهيت من صور القبل")]],
                resize_keyboard=True,
                one_time_keyboard=False
            )
        )
    else:
        await update.message.reply_text(
            f"✅ *{loc_data['name']}*\n\n"
            "📸 أرسل الصور الآن — واحدة أو أكثر\n"
            "⬇️ عند الانتهاء اضغط الزر أدناه:",
            parse_mode="Markdown",
            reply_markup=ReplyKeyboardMarkup(
                [[KeyboardButton("✅ انتهيت من الصور")]],
                resize_keyboard=True,
                one_time_keyboard=False
            )
        )
    return True



async def next_after_location(session, update=None):
    """بعد تسجيل الموقع — توعية تسأل عن النوع، تفتيش تسأل عن الملاحظة"""
    if session.get("dept") == "awareness":
        if update:
            await update.message.reply_text(
                "📋 ما نوع التوعية؟",
                reply_markup=ReplyKeyboardRemove()
            )
        return WAITING_AWARENESS_TYPE
    else:
        if update:
            kb = [[KeyboardButton("⏭️ بدون ملاحظة")]]
            await update.message.reply_text(
                "📝 ملاحظة للصورة؟",
                reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
            )
        return WAITING_NOTE


# ══════════════════════════════════════════════════
#  نوع التوعية (قسم التوعية فقط)
# ══════════════════════════════════════════════════

async def finish_collecting(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """الضغط على انتهيت من الصور — يحفظ الصور ويكمل"""
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})

    cmd = await handle_commands(update, context, COLLECTING_PHOTOS)
    if cmd is not None:
        return cmd

    photos = session.get("collected_photos", [])
    if not photos:
        await update.message.reply_text(
            "❌ لم ترسل أي صورة! أرسل صورة أولاً.",
            reply_markup=ReplyKeyboardMarkup(
                [[KeyboardButton("✅ انتهيت من الصور")]],
                resize_keyboard=True
            )
        )
        return COLLECTING_PHOTOS

    loc    = session["last_location"]
    coords = session.get("last_coords_dms", "")

    is_ba = session.get("report_type") == "before_after"

    for ph in photos:
        session["photos"].append({
            "photo":          ph,
            "photo_after":    None,   # فارغ — يُعبأ لاحقاً
            "location":       loc,
            "coords_dms":     coords,
            "plot_no":        "",
            "note":           "",
            "awareness_type": "",
            "type":           "before_after" if is_ba else "normal"
        })

    count_added = len(photos)
    total       = len(session["photos"])
    session["collecting"]        = False
    session["collected_photos"]  = []
    session["collecting_msg_id"] = None
    session["last_batch_count"]  = count_added

    await update.message.reply_text(
        f"✅ تم إضافة *{count_added}* صورة — إجمالي: *{total}*",
        parse_mode="Markdown",
        reply_markup=ReplyKeyboardRemove()
    )

    if is_ba:
        # في وضع قبل/بعد: اسأل عن ملاحظة ثم اعرض خيار إضافة موقع آخر أو إنشاء التقرير
        kb = [[KeyboardButton("⏭️ بدون ملاحظة")]]
        await update.message.reply_text(
            "📝 ملاحظة لهذا الموقع؟",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_BA_NOTE
    return await next_after_location(session, update)


async def receive_awareness_type(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    text    = update.message.text.strip()

    cmd = await handle_commands(update, context, WAITING_AWARENESS_TYPE)
    if cmd is not None:
        return cmd

    if session["photos"]:
        # عدد الصور التي أُضيفت في آخر دفعة
        last_batch = session.get("last_batch_count", 1)
        # تطبيق نوع التوعية على كل صور الدفعة الأخيرة
        for photo in session["photos"][-last_batch:]:
            photo["awareness_type"] = text

    kb = [[KeyboardButton("⏭️ بدون ملاحظة")]]
    await update.message.reply_text(
        f"✅ نوع التوعية: *{text}*\n\n📝 ملاحظات إضافية؟",
        parse_mode="Markdown",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )
    return WAITING_NOTE


# ══════════════════════════════════════════════════
#  الملاحظة
# ══════════════════════════════════════════════════

async def receive_note(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    text    = update.message.text.strip()

    cmd = await handle_commands(update, context, WAITING_NOTE)
    if cmd is not None:
        return cmd

    note = "" if "بدون ملاحظة" in text else text
    if session["photos"]:
        session["photos"][-1]["note"] = note

    # سؤال القسيمة للتفتيش فقط
    if session.get("dept") != "awareness":
        kb = [[KeyboardButton("⏭️ بدون قسيمة")]]
        await update.message.reply_text(
            "🏠 رقم القسيمة / العنوان؟",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_PLOT

    count = len(session["photos"])
    kb = [[KeyboardButton("📸 إضافة صورة")], [KeyboardButton("✅ إنشاء التقرير")]]
    await update.message.reply_text(
        f"✅ الصورة {count} مضافة\n\nإضافة المزيد أو إنشاء التقرير؟",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )
    return ADD_MORE


async def receive_plot(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """\u0627\u0633\u062a\u0642\u0628\u0627\u0644 \u0631\u0642\u0645 \u0627\u0644\u0642\u0633\u064a\u0645\u0629 / \u0627\u0644\u0639\u0646\u0648\u0627\u0646"""
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    text    = update.message.text.strip()

    cmd = await handle_commands(update, context, WAITING_PLOT)
    if cmd is not None:
        return cmd

    plot = "" if "\u0628\u062f\u0648\u0646 \u0642\u0633\u064a\u0645\u0629" in text else text
    if session["photos"]:
        session["photos"][-1]["plot_no"] = plot

    count = len(session["photos"])
    kb = [[KeyboardButton("\U0001f4f8 \u0625\u0636\u0627\u0641\u0629 \u0635\u0648\u0631\u0629")], [KeyboardButton("\u2705 \u0625\u0646\u0634\u0627\u0621 \u0627\u0644\u062a\u0642\u0631\u064a\u0631")]]
    await update.message.reply_text(
        f"\u2705 \u0627\u0644\u0635\u0648\u0631\u0629 {count} \u0645\u0636\u0627\u0641\u0629\n\n\u0625\u0636\u0627\u0641\u0629 \u0627\u0644\u0645\u0632\u064a\u062f \u0623\u0648 \u0625\u0646\u0634\u0627\u0621 \u0627\u0644\u062a\u0642\u0631\u064a\u0631\u061f",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )
    return ADD_MORE


async def add_more(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})
    text    = update.message.text

    cmd = await handle_commands(update, context, ADD_MORE)
    if cmd is not None:
        return cmd

    if "إضافة صورة" in text or "📸" in text or "موقع جديد" in text:
        kb = [[KeyboardButton("📍 مشاركة موقعي", request_location=True)]]
        await update.message.reply_text(
            f"📍 أرسل الموقع الجديد:\n"
            f"(التقرير يحتوي على {len(session.get('photos', []))} صورة حتى الآن)",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_LOCATION
    elif "إنشاء التقرير" in text or "✅" in text:
        return await do_send(update, context)
    else:
        reply = ai_reply(text, session)
        kb = [[KeyboardButton("📸 إضافة صورة")], [KeyboardButton("✅ إنشاء التقرير")]]
        await update.message.reply_text(reply, reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True))
        return ADD_MORE


# ══════════════════════════════════════════════════
#  تقرير قبل وبعد (تفتيش فقط)
# ══════════════════════════════════════════════════

async def receive_before_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    photo   = update.message.photo[-1]
    file    = await context.bot.get_file(photo.file_id)
    pb      = await file.download_as_bytearray()
    session["before_photo"] = bytes(pb)
    await update.message.reply_text(
        "✅ صورة 'قبل' مستلمة\n\n📸 أرسل صورة *بعد*:",
        parse_mode="Markdown", reply_markup=ReplyKeyboardRemove()
    )
    return WAITING_AFTER_PHOTO


async def receive_after_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    photo   = update.message.photo[-1]
    file    = await context.bot.get_file(photo.file_id)
    pb      = await file.download_as_bytearray()
    session["after_photo"] = bytes(pb)
    await ask_location(update)
    return WAITING_BA_LOC


async def receive_ba_location_gps(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    loc     = update.message.location
    await process_ba_location(update, session, loc.latitude, loc.longitude)
    return WAITING_BA_NOTE


async def receive_ba_location_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    lat, lng = parse_coords(update.message.text.strip())
    if lat:
        await process_ba_location(update, session, lat, lng)
        return WAITING_BA_NOTE
    await update.message.reply_text("❌ أرسل رابط أو إحداثيات صحيحة")
    return WAITING_BA_LOC


async def process_ba_location(update, session, lat, lng):
    await update.message.reply_text("⏳ جاري تحديد الموقع...")
    loc_data = get_address(lat, lng)
    if not loc_data["in_liwa"]:
        await update.message.reply_text("⚠️ الموقع خارج نطاق ليوا! أعد الإرسال.")
        await ask_location(update)
        return
    session["photos"].append({
        "photo":       session["before_photo"],
        "photo_after": session["after_photo"],
        "location":    loc_data,
        "note":        "",
        "type":        "before_after"
    })
    session["last_location"] = loc_data
    kb = [[KeyboardButton("⏭️ بدون ملاحظة")]]
    await update.message.reply_text(
        f"✅ *{loc_data['name']}*\n\n📝 ملاحظة؟",
        parse_mode="Markdown",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )


async def receive_ba_note(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    session = user_sessions[user_id]
    text    = update.message.text.strip()
    if session["photos"]:
        session["photos"][-1]["note"] = "" if "بدون ملاحظة" in text else text
    kb = [[KeyboardButton("📸 إضافة موقع آخر")], [KeyboardButton("✅ إنشاء التقرير")]]
    await update.message.reply_text(
        "✅ تمت الإضافة\n\nإضافة موقع آخر أو إنشاء التقرير؟",
        reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
    )
    return ADD_MORE


# ══════════════════════════════════════════════════
#  إنشاء PowerPoint
# ══════════════════════════════════════════════════

def create_pptx(session: dict) -> BytesIO:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    is_awareness = session.get("dept") == "awareness"
    name         = session.get("name", "")
    date         = session.get("date", "")
    now_time     = datetime.now().strftime("%H:%M")
    photos       = session.get("photos", [])
    header_color = BROWN
    report_title = "تقرير توعية" if is_awareness else "تقرير رصد"
    label_type   = "نوع التوعية" if is_awareness else "نوع الرصد"

    def add_bg(slide):
        s = slide.shapes.add_shape(1, 0, 0, W, H)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xF5, 0xF3, 0xF0)
        s.line.fill.background()

    def add_corners(slide):
        sz = Inches(0.35)
        for x, y in [(0,0),(W-sz,0),(0,H-sz),(W-sz,H-sz)]:
            s = slide.shapes.add_shape(1, x, y, sz, sz)
            s.fill.solid(); s.fill.fore_color.rgb = header_color
            s.line.fill.background()

    def txt(slide, text, x, y, w, h, size=10, bold=False,
            color=BLACK, align=PP_ALIGN.RIGHT, italic=False):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = color
        run.font.name = "Arial"

    def photo_card(slide, photo_bytes, x, y, w, h, entry):
        """
        صورة مربعة/عمودية نسبياً مع بيانات تحتها:
        الموقع، الإحداثيات، نوع التوعية/الرصد، الملاحظة
        """
        loc_name      = entry.get("location", {}).get("name", "")
        note          = entry.get("note", "")
        awareness_type= entry.get("awareness_type", "")
        coords_dms    = entry.get("coords_dms", "")
        plot_no       = entry.get("plot_no", "")

        # بيانات تُعرض تحت الصورة
        items = []
        if loc_name:       items.append(("الموقع", loc_name))
        if coords_dms:     items.append(("الإحداثيات", coords_dms))
        if plot_no:        items.append(("القسيمة", plot_no))
        if awareness_type: items.append((label_type, awareness_type))
        if note:           items.append(("الملاحظة", note))

        # ارتفاع ثابت للبيانات: 0.42 بوصة لكل عنصر
        ROW_H  = Inches(0.42)
        info_h = ROW_H * max(len(items), 1) + Inches(0.1)
        # الصورة تأخذ ما تبقى — نحافظ على نسبة 4:3 تقريباً
        img_h  = h - info_h - Inches(0.08)
        pad    = Inches(0.08)

        # إطار الكارد
        frame = slide.shapes.add_shape(1, x, y, w, h)
        frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = header_color; frame.line.width = Pt(1.2)

        # الصورة — contain كامل (نسبة محفوظة، توسيط)
        try:
            img_stream = BytesIO(photo_bytes)
            pil_img = PILImage.open(BytesIO(photo_bytes))
            iw, ih  = pil_img.size
            ratio   = iw / ih

            max_w = w - pad * 2
            max_h = img_h - pad * 2

            if (max_w / ratio) <= max_h:
                pic_w = max_w
                pic_h = max_w / ratio
            else:
                pic_h = max_h
                pic_w = max_h * ratio

            pic_x = x + (w - pic_w) / 2
            pic_y = y + pad + (max_h - pic_h) / 2
            slide.shapes.add_picture(img_stream, pic_x, pic_y, pic_w, pic_h)
        except Exception as e:
            logger.warning(f"Photo: {e}")

        # منطقة البيانات
        iy = y + img_h
        if items:
            # خلفية البيانات
            bg = slide.shapes.add_shape(1, x, iy, w, info_h + Inches(0.08))
            bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF2, 0xEE, 0xE8)
            bg.line.fill.background()

            # خط ملون في الأعلى
            ac = slide.shapes.add_shape(1, x, iy, w, Inches(0.035))
            ac.fill.solid(); ac.fill.fore_color.rgb = header_color
            ac.line.fill.background()

            for i, (label, val) in enumerate(items):
                by = iy + Inches(0.05) + i * ROW_H
                # تسمية (صغيرة ملونة)
                txt(slide, label + ":",
                    x + Inches(0.1), by, w - Inches(0.2), Inches(0.18),
                    size=7, bold=True, color=header_color, align=PP_ALIGN.RIGHT)
                # قيمة
                txt(slide, val,
                    x + Inches(0.1), by + Inches(0.18), w - Inches(0.2), Inches(0.22),
                    size=8.5, color=BLACK, align=PP_ALIGN.RIGHT)

    # ── الغلاف
    def make_cover():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide); add_corners(slide)

        # شريط علوي
        bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.15))
        bar.fill.solid(); bar.fill.fore_color.rgb = header_color
        bar.line.fill.background()

        # عنوان
        txt(slide, report_title,
            Inches(1.5), Inches(2.0), W-Inches(3), Inches(1.2),
            size=40, bold=True, color=header_color, align=PP_ALIGN.CENTER)

        # خط فاصل
        line = slide.shapes.add_shape(1, Inches(3), Inches(3.4), W-Inches(6), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = header_color
        line.line.fill.background()

        # بيانات الغلاف — بدون "الجهة"
        lbl_inspector = "المسؤول:" if is_awareness else "المراقب:"
        lbl_count     = "عدد صور التوعية:" if is_awareness else "عدد عمليات الرصد:"
        info = [
            (lbl_inspector, name),
            ("التاريخ:", date),
            ("الوقت:", now_time),
            (lbl_count, str(len(photos))),
        ]
        for i, (lbl, val) in enumerate(info):
            iy = Inches(3.65) + i * Inches(0.52)
            txt(slide, lbl, Inches(7.5), iy, Inches(2.0), Inches(0.42),
                size=13, bold=True, color=header_color, align=PP_ALIGN.RIGHT)
            txt(slide, val, Inches(3.5), iy, Inches(4.0), Inches(0.42),
                size=13, color=BLACK, align=PP_ALIGN.RIGHT)

        # شريط سفلي
        bbar = slide.shapes.add_shape(1, 0, H-Inches(0.15), W, Inches(0.15))
        bbar.fill.solid(); bbar.fill.fore_color.rgb = header_color
        bbar.line.fill.background()
        txt(slide, f"{date}  {now_time}",
            Inches(0.4), H-Inches(0.42), Inches(4), Inches(0.3),
            size=9, color=WHITE, align=PP_ALIGN.LEFT)

    # ── صفحة الصور (شبكة مرنة 1-4 صور)
    PHOTOS_PER_PAGE = 4

    def make_grid_slide(entries, page_num):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide); add_corners(slide)

        # شريط علوي
        hbar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.15))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = header_color
        hbar.line.fill.background()

        txt(slide, f"صفحة {page_num}",
            Inches(0.5), Inches(0.2), Inches(4), Inches(0.3),
            size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        txt(slide, f"{date}  {now_time}",
            W-Inches(4.5), Inches(0.2), Inches(4), Inches(0.3),
            size=9, color=WHITE, align=PP_ALIGN.LEFT)

        n     = len(entries)
        pad   = Inches(0.22)
        gap   = Inches(0.15)
        top_y = Inches(0.55)
        bot_y = H - Inches(0.15)

        # تخطيط الشبكة
        if n == 1:
            cols, rows = 1, 1
        elif n == 2:
            cols, rows = 2, 1
        elif n == 3:
            cols, rows = 3, 1
        else:  # 4
            cols, rows = 2, 2

        avail_w = W - pad*2 - gap*(cols-1)
        avail_h = bot_y - top_y - gap*(rows-1)
        card_w  = avail_w / cols
        card_h  = avail_h / rows

        for idx, entry in enumerate(entries):
            col = idx % cols
            row = idx // cols
            cx  = pad + col*(card_w + gap)
            cy  = top_y + row*(card_h + gap)
            photo_card(slide, entry["photo"], cx, cy, card_w, card_h, entry)

        # شريط سفلي
        bbar = slide.shapes.add_shape(1, 0, H-Inches(0.15), W, Inches(0.15))
        bbar.fill.solid(); bbar.fill.fore_color.rgb = header_color
        bbar.line.fill.background()

    # ── صفحة قبل/بعد
    def draw_ba_cell(slide, entry, cx, cy, cw, ch, seq_num):
        """
        خلية قبل/بعد:
        - يمين: صورة القبل
        - يسار: مربع فارغ للبعد
        - بيانات في الوسط أسفل كل من القبل والبعد
        """
        p_gap      = Inches(0.12)
        DATA_H     = Inches(0.72)
        img_area_h = ch - DATA_H - Inches(0.04)
        half_w     = (cw - p_gap) / 2

        # ── القبل (يمين الخلية)
        bx = cx + half_w + p_gap   # يمين = الجزء الأيمن
        ax = cx                    # يسار = الجزء الأيسر

        # رقم تسلسلي في أعلى المنتصف
        num_cx = cx + (cw - Inches(0.32)) / 2
        num_box = slide.shapes.add_shape(1, num_cx, cy, Inches(0.32), Inches(0.22))
        num_box.fill.solid(); num_box.fill.fore_color.rgb = header_color
        num_box.line.fill.background()
        txt(slide, str(seq_num), num_cx, cy, Inches(0.32), Inches(0.22),
            size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # إطار القبل
        frame_b = slide.shapes.add_shape(1, bx, cy, half_w, img_area_h)
        frame_b.fill.solid(); frame_b.fill.fore_color.rgb = WHITE
        frame_b.line.color.rgb = header_color; frame_b.line.width = Pt(1.2)

        txt(slide, "◄ قبل", bx + Inches(0.05), cy + Inches(0.03),
            half_w - Inches(0.1), Inches(0.18),
            size=7, bold=True, color=header_color, align=PP_ALIGN.RIGHT)

        # صورة القبل (contain)
        try:
            img_bytes = entry["photo"]
            pil_img   = PILImage.open(BytesIO(img_bytes))
            iw, ih    = pil_img.size
            ratio     = iw / ih
            p_pad     = Inches(0.06)
            max_w     = half_w - p_pad * 2
            max_h     = img_area_h - p_pad * 2
            if (max_w / ratio) <= max_h:
                pic_w = max_w; pic_h = max_w / ratio
            else:
                pic_h = max_h; pic_w = max_h * ratio
            slide.shapes.add_picture(
                BytesIO(img_bytes),
                bx + (half_w - pic_w) / 2,
                cy + p_pad + (max_h - pic_h) / 2,
                pic_w, pic_h
            )
        except Exception as e:
            logger.warning(f"BA photo: {e}")

        # إطار البعد الفارغ (يسار)
        frame_a = slide.shapes.add_shape(1, ax, cy, half_w, img_area_h)
        frame_a.fill.solid(); frame_a.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        frame_a.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        frame_a.line.width = Pt(1.2)
        frame_a.line.dash_style = 4

        txt(slide, "بعد ►", ax + Inches(0.05), cy + Inches(0.03),
            half_w - Inches(0.1), Inches(0.18),
            size=7, bold=True, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)

        txt(slide, "أضف صورة البعد هنا",
            ax, cy + img_area_h / 2 - Inches(0.18), half_w, Inches(0.36),
            size=8, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

        # ── البيانات في الوسط - شريط مشترك تحت الصورتين
        dy  = cy + img_area_h + Inches(0.04)
        bg  = slide.shapes.add_shape(1, cx, dy, cw, DATA_H)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF2, 0xEE, 0xE8)
        bg.line.fill.background()

        ac = slide.shapes.add_shape(1, cx, dy, cw, Inches(0.03))
        ac.fill.solid(); ac.fill.fore_color.rgb = header_color
        ac.line.fill.background()

        loc   = entry.get("location", {}).get("name", "")
        coord = entry.get("coords_dms", "")
        note  = entry.get("note", "")

        rows = []
        if loc:   rows.append(("الموقع:", loc))
        if coord: rows.append(("الإحداثيات:", coord))
        if note:  rows.append(("ملاحظة:", note))

        ROW = Inches(0.33)
        for ri, (lb, vl) in enumerate(rows[:2]):
            ry = dy + Inches(0.04) + ri * ROW
            txt(slide, lb,
                cx + Inches(0.1), ry, cw * 0.35, Inches(0.15),
                size=6.5, bold=True, color=header_color, align=PP_ALIGN.RIGHT)
            txt(slide, vl,
                cx + Inches(0.1), ry + Inches(0.15), cw - Inches(0.2), Inches(0.18),
                size=8, color=BLACK, align=PP_ALIGN.RIGHT)

    def make_ba_page(entries, page_num, seq_start):
        """
        صفحة تحتوي على 2-4 خلايا قبل/بعد.
        التخطيط:
          2 صور → صفّان (1 عمود)
          3 صور → صفّان + 1 (الثالثة وحدها في صف)
          4 صور → صفّان × عمودان
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide); add_corners(slide)

        # شريط علوي
        hbar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.45))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = header_color
        hbar.line.fill.background()

        txt(slide, "قبل وبعد الرصد",
            Inches(0.4), Inches(0.08), W * 0.5, Inches(0.3),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
        txt(slide, f"صفحة {page_num}  |  {date}  {now_time}",
            W * 0.5, Inches(0.08), W * 0.5 - Inches(0.4), Inches(0.3),
            size=9, color=WHITE, align=PP_ALIGN.LEFT)

        n = len(entries)
        pad  = Inches(0.18)
        gap  = Inches(0.14)
        top  = Inches(0.5)
        bot  = H - Inches(0.18)

        # تحديد عدد الأعمدة والصفوف
        if n <= 2:
            cols, rows_count = 1, n
        elif n == 3:
            cols, rows_count = 1, 3
        else:  # 4
            cols, rows_count = 2, 2

        avail_w = W - pad * 2 - gap * (cols - 1)
        avail_h = bot - top - gap * (rows_count - 1)
        cell_w  = avail_w / cols
        cell_h  = avail_h / rows_count

        for idx, entry in enumerate(entries):
            col = idx % cols
            row = idx // cols
            cx  = pad + col * (cell_w + gap)
            cy  = top + row * (cell_h + gap)
            draw_ba_cell(slide, entry, cx, cy, cell_w, cell_h, seq_start + idx)

        # شريط سفلي
        bbar = slide.shapes.add_shape(1, 0, H - Inches(0.15), W, Inches(0.15))
        bbar.fill.solid(); bbar.fill.fore_color.rgb = header_color
        bbar.line.fill.background()

    def make_closing():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide); add_corners(slide)
        fw, fh = Inches(5.5), Inches(2.8)
        fx, fy = (W-fw)/2, (H-fh)/2
        frame = slide.shapes.add_shape(1, fx, fy, fw, fh)
        frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = header_color; frame.line.width = Pt(2)
        txt(slide, "شكـراً", fx, fy+Inches(0.35), fw, Inches(1.4),
            size=44, bold=True, color=header_color, align=PP_ALIGN.CENTER)
        txt(slide, "Thank You", fx, fy+Inches(1.75), fw, Inches(0.6),
            size=16, color=RGBColor(0xAA,0x66,0x44),
            align=PP_ALIGN.CENTER, italic=True)
        bbar = slide.shapes.add_shape(1, 0, H-Inches(0.15), W, Inches(0.15))
        bbar.fill.solid(); bbar.fill.fore_color.rgb = header_color
        bbar.line.fill.background()
        txt(slide, f"{date}  {now_time}",
            Inches(0.4), H-Inches(0.42), Inches(4), Inches(0.3),
            size=9, color=WHITE, align=PP_ALIGN.LEFT)

    # ── بناء الشرائح
    make_cover()

    normal  = [e for e in photos if e.get("type") != "before_after"]
    ba_list = [e for e in photos if e.get("type") == "before_after"]

    page = 1
    for i in range(0, len(normal), PHOTOS_PER_PAGE):
        make_grid_slide(normal[i:i+PHOTOS_PER_PAGE], page)
        page += 1

    def ba_distribute(entries, min_per=2, max_per=4):
        """توزيع صور قبل/بعد على صفحات: حد أدنى 2، حد أقصى 4"""
        n = len(entries)
        if n == 0: return []
        if n <= max_per: return [entries]
        pages = []
        remaining = list(entries)
        while remaining:
            if len(remaining) <= max_per:
                if len(remaining) < min_per and pages:
                    # أعد صورة من الصفحة السابقة لتكملة الحد الأدنى
                    moved = pages[-1].pop()
                    remaining.insert(0, moved)
                pages.append(remaining[:])
                remaining = []
            else:
                pages.append(remaining[:max_per])
                remaining = remaining[max_per:]
        return pages

    seq = 1
    for page_ba, chunk in enumerate(ba_distribute(ba_list), 1):
        make_ba_page(chunk, page_ba, seq)
        seq += len(chunk)

    make_closing()

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ══════════════════════════════════════════════════
#  إرسال الإيميل
# ══════════════════════════════════════════════════

def send_email_to(buf: BytesIO, session: dict, fname: str, to_email: str) -> bool:
    """إرسال التقرير لإيميل محدد"""
    if not RESEND_API_KEY:
        return False
    try:
        is_awareness = session.get("dept") == "awareness"
        dept_ar      = "توعية" if is_awareness else "تفتيش"
        locs = list(dict.fromkeys([
            p["location"]["name"] for p in session.get("photos",[]) if p.get("location")
        ]))
        body = (
            f"تقرير {dept_ar} جديد — {BOT_NAME} {BOT_VERSION}\n\n"
            f"المسؤول: {session.get('name','—')}\n"
            f"التاريخ: {session.get('date','—')}\n"
            f"الصور: {len(session.get('photos',[]))}\n"
            f"المواقع: {', '.join(locs) or '—'}\n"
        )
        file_content = base64.b64encode(buf.read()).decode()
        response = req.post(
            "https://api.resend.com/emails",
            headers={"Authorization": f"Bearer {RESEND_API_KEY}", "Content-Type": "application/json"},
            json={
                "from":    f"بوت ليوا <onboarding@resend.dev>",
                "to":      [to_email],
                "subject": f"تقرير {dept_ar} — {session.get('name','')} — {session.get('date','')}",
                "text":    body,
                "attachments": [{"filename": fname, "content": file_content}]
            },
            timeout=15
        )
        return response.status_code == 200
    except Exception as e:
        logger.error(f"Email: {e}")
        return False


def send_email(buf: BytesIO, session: dict, fname: str) -> bool:
    if not RESEND_API_KEY:
        return False
    try:
        is_awareness = session.get("dept") == "awareness"
        dept_ar      = "توعية" if is_awareness else "تفتيش"

        locs = list(dict.fromkeys([
            p["location"]["name"] for p in session.get("photos",[]) if p.get("location")
        ]))
        body = (
            f"تقرير {dept_ar} جديد — {BOT_NAME} {BOT_VERSION}\n\n"
            f"المسؤول: {session.get('name','—')}\n"
            f"التاريخ: {session.get('date','—')}\n"
            f"الصور: {len(session.get('photos',[]))}\n"
            f"المواقع: {', '.join(locs) or '—'}\n"
        )

        # تحويل الملف لـ base64
        file_content = base64.b64encode(buf.read()).decode()

        # إرسال عبر Resend API
        response = req.post(
            "https://api.resend.com/emails",
            headers={
                "Authorization": f"Bearer {RESEND_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "from":    "بوت ليوا <onboarding@resend.dev>",
                "to":      [ADMIN_EMAIL],
                "subject": f"تقرير {dept_ar} — {session.get('name','')} — {session.get('date','')}",
                "text":    body,
                "attachments": [{
                    "filename": fname,
                    "content":  file_content
                }]
            },
            timeout=15
        )

        if response.status_code == 200:
            return True
        else:
            logger.error(f"Resend error: {response.text}")
            return False

    except Exception as e:
        logger.error(f"Email: {e}")
        return False


# ══════════════════════════════════════════════════
#  إنشاء وإرسال التقرير
# ══════════════════════════════════════════════════

async def do_send(update, context):
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})

    if not session.get("photos"):
        await update.message.reply_text("❌ لا توجد صور!", reply_markup=ReplyKeyboardRemove())
        return ConversationHandler.END

    await update.message.reply_text("⏳ جاري مراجعة البيانات وإنشاء التقرير...", reply_markup=ReplyKeyboardRemove())

    try:
        is_awareness = session.get("dept") == "awareness"
        dept_ar      = "توعية" if is_awareness else "رصد"
        prepared_session = agent_prepare_session(session)
        buf   = create_pptx(prepared_session)
        rand3 = random.randint(100, 999)
        fname = f"2025{rand3}_{dept_ar}.pptx"

        buf.seek(0)
        await update.message.reply_document(
            document=buf, filename=fname,
            caption=(
                f"✅ تقرير {dept_ar} جاهز — {BOT_VERSION}\n"
                f"👤 {session['name']}\n"
                f"📸 {len(session['photos'])} صورة\n"
                f"📅 {session['date']}"
            )
        )

        # تحديث الإحصائيات
        report_stats["total"] += 1
        if is_awareness:
            report_stats["awareness"] += 1
        else:
            report_stats["inspection"] += 1

        # حفظ الملف في session للإرسال لاحقاً
        buf.seek(0)
        session["report_buf"]   = buf.read()
        session["report_fname"] = fname

        # سؤال الإيميل
        kb = [[KeyboardButton("📧 نعم، أرسله")], [KeyboardButton("❌ لا شكراً")]]
        await update.message.reply_text(
            "📨 هل تود إرسال التقرير عبر البريد الإلكتروني؟",
            reply_markup=ReplyKeyboardMarkup(kb, resize_keyboard=True)
        )
        return WAITING_EMAIL

    except Exception as e:
        logger.error(f"Send error: {e}")
        await update.message.reply_text(f"❌ خطأ: {e}\n\nاكتب /start وحاول مجدداً")
        user_sessions.pop(user_id, None)
        return ConversationHandler.END


async def receive_email_choice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """استقبال قرار الإرسال بالإيميل"""
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})
    text    = update.message.text.strip()
    name    = session.get("name", "")

    if "نعم" in text or "📧" in text:
        await update.message.reply_text(
            "📧 اكتب البريد الإلكتروني المراد الإرسال إليه:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_EMAIL

    # لا — إنهاء بدون إرسال
    await update.message.reply_text(
        f"✅ شكراً لك أخ {name}!\n🎉 تقرير رقم {report_stats['total']} مكتمل.\n\nاكتب /start لتقرير جديد",
        reply_markup=ReplyKeyboardRemove()
    )
    user_sessions.pop(user_id, None)
    return ConversationHandler.END


async def receive_email_address(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """استقبال البريد الإلكتروني وإرسال التقرير"""
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})
    email   = update.message.text.strip()
    name    = session.get("name", "")

    # التحقق من صحة الإيميل
    import re as _re
    if not _re.match(r"[^@]+@[^@]+\.[^@]+", email):
        await update.message.reply_text(
            "❌ البريد الإلكتروني غير صحيح، أعد المحاولة:\n(مثال: name@gmail.com)"
        )
        return WAITING_EMAIL

    await update.message.reply_text(f"📤 جاري الإرسال إلى {email}...")

    try:
        buf   = BytesIO(session.get("report_buf", b""))
        fname = session.get("report_fname", "report.pptx")

        # إرسال للإيميل المحدد
        ok = send_email_to(buf, session, fname, email)
        if ok:
            await update.message.reply_text(
                f"✅ تم إرسال التقرير إلى:\n📩 {email}"
            )
        else:
            await update.message.reply_text("⚠️ فشل الإرسال — تأكد من صحة الإيميل وحاول مجدداً.")
    except Exception as e:
        logger.error(f"Email error: {e}")
        await update.message.reply_text("⚠️ حدث خطأ أثناء الإرسال.")

    await update.message.reply_text(
        f"شكراً لك أخ {name}! 🎉\nاكتب /start لتقرير جديد",
        reply_markup=ReplyKeyboardRemove()
    )
    user_sessions.pop(user_id, None)
    return ConversationHandler.END


async def receive_email_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """يتعامل مع مرحلة WAITING_EMAIL:
    إذا اختار نعم في المرة الأولى → يطلب الإيميل
    إذا أرسل إيميل → يرسل التقرير
    إذا اختار لا → ينهي
    """
    user_id = update.effective_user.id
    session = user_sessions.get(user_id, {})
    text    = update.message.text.strip()
    name    = session.get("name", "")

    # لا — إنهاء بدون إرسال
    if "لا" in text or "❌" in text:
        await update.message.reply_text(
            f"شكراً لك أخ {name}! 🎉\n"
            f"🎊 تقرير رقم {report_stats['total']} مكتمل.\n\n"
            "اكتب /start لتقرير جديد",
            reply_markup=ReplyKeyboardRemove()
        )
        user_sessions.pop(user_id, None)
        return ConversationHandler.END

    # نعم — طلب الإيميل
    if "نعم" in text or "📧" in text:
        session["awaiting_email_input"] = True
        await update.message.reply_text(
            "📧 اكتب البريد الإلكتروني:",
            reply_markup=ReplyKeyboardRemove()
        )
        return WAITING_EMAIL

    # محاولة معالجة النص كإيميل
    import re as _re
    if session.get("awaiting_email_input") and _re.match(r"[^@]+@[^@]+\.[^@]+", text):
        await update.message.reply_text(f"📤 جاري الإرسال إلى {text}...")
        try:
            buf   = BytesIO(session.get("report_buf", b""))
            fname = session.get("report_fname", "report.pptx")
            ok    = send_email_to(buf, session, fname, text)
            if ok:
                await update.message.reply_text(f"✅ تم الإرسال بنجاح إلى:\n📩 {text}")
            else:
                await update.message.reply_text("⚠️ فشل الإرسال — تأكد من صحة الإيميل.")
        except Exception as e:
            logger.error(f"Email send: {e}")
            await update.message.reply_text("⚠️ حدث خطأ أثناء الإرسال.")

        await update.message.reply_text(
            f"شكراً لك أخ {name}! 🎉\nاكتب /start لتقرير جديد",
            reply_markup=ReplyKeyboardRemove()
        )
        user_sessions.pop(user_id, None)
        return ConversationHandler.END

    # إيميل غير صحيح
    if session.get("awaiting_email_input"):
        await update.message.reply_text(
            "❌ البريد غير صحيح، أعد الكتابة:\n(مثال: name@gmail.com)"
        )
        return WAITING_EMAIL

    return WAITING_EMAIL


# ══════════════════════════════════════════════════
#  تشغيل البوت
# ══════════════════════════════════════════════════

def main():
    app = Application.builder().token(TELEGRAM_TOKEN).build()

    conv = ConversationHandler(
        entry_points=[
            CommandHandler("start", start),
        ],
        states={
            WAITING_CODE:        [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_code)],
            WAITING_NAME:        [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_name)],
            MAIN_MENU:           [MessageHandler(filters.TEXT & ~filters.COMMAND, main_menu)],
            WAITING_PHOTO:       [MessageHandler(filters.PHOTO, receive_photo)],
            WAITING_LOCATION:    [
                MessageHandler(filters.LOCATION, receive_location_gps),
                MessageHandler(filters.TEXT & ~filters.COMMAND, receive_location_text),
            ],
            WAITING_AWARENESS_TYPE: [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_awareness_type)],
            WAITING_NOTE:        [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_note)],
            WAITING_PLOT:        [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_plot)],
            ADD_MORE:            [
                MessageHandler(filters.PHOTO, receive_photo),
                MessageHandler(filters.TEXT & ~filters.COMMAND, add_more),
            ],
            COLLECTING_PHOTOS:   [
                MessageHandler(filters.PHOTO, receive_photo),
                MessageHandler(filters.TEXT & ~filters.COMMAND, finish_collecting),
            ],
            WAITING_BEFORE_PHOTO: [MessageHandler(filters.PHOTO, receive_before_photo)],
            WAITING_AFTER_PHOTO:  [MessageHandler(filters.PHOTO, receive_after_photo)],
            WAITING_BA_LOC:      [
                MessageHandler(filters.LOCATION, receive_ba_location_gps),
                MessageHandler(filters.TEXT & ~filters.COMMAND, receive_ba_location_text),
            ],
            WAITING_BA_NOTE:     [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_ba_note)],
            WAITING_EMAIL:       [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_email_handler)],
        },
        fallbacks=[CommandHandler("start", start)],
        allow_reentry=True,
    )

    app.add_handler(conv)
    print(f"🤖 {BOT_NAME} {BOT_VERSION} يعمل...")
    app.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
